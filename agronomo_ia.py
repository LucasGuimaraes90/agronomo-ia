import streamlit as st
import anthropic
import base64
import io
import re
from datetime import datetime
import json

# в”Ђв”Ђ PDF в”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђ
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.lib import colors
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table,
    TableStyle, HRFlowable, KeepTogether
)

# в”Ђв”Ђ Excel в”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђ
import openpyxl
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side, numbers

# в”Ђв”Ђ Word в”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђ
from docx import Document
from docx.shared import Pt, RGBColor, Cm, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

# в”Ђв”Ђ PowerPoint в”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђ
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt, Emu
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN

# в•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђ
# CONFIG
# в•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђ
st.set_page_config(
    page_title="AgrГґnomo IA вЂ” Guara Agro",
    page_icon="рџЊ±",
    layout="centered",
    initial_sidebar_state="collapsed",
)

# в”Ђв”Ђ CSS tema verde claro / white в”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђ
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;600;700;900&display=swap');

html, body, [class*="css"] { font-family: 'Inter', sans-serif; }

/* fundo geral вЂ” branco */
.stApp { background: #ffffff; color: #1a2e1a; }

/* sidebar вЂ” verde muito suave */
section[data-testid="stSidebar"] { background: #f1f8f1; border-right: 1px solid #c8e6c9; }

/* chat bubbles вЂ” usuГЎrio */
[data-testid="stChatMessage"][data-testid*="user"] {
    background: #e8f5e9;
    border-radius: 12px;
    margin: 6px 0;
    border: 1px solid #c8e6c9;
}

/* chat bubbles вЂ” assistente */
[data-testid="stChatMessage"] {
    background: #f9fdf9;
    border-radius: 12px;
    margin: 6px 0;
    border: 1px solid #e0f0e0;
}

/* texto dentro dos bubbles */
[data-testid="stChatMessage"] p,
[data-testid="stChatMessage"] li,
[data-testid="stChatMessage"] td,
[data-testid="stChatMessage"] th { color: #1a2e1a !important; }

/* tabelas no chat */
[data-testid="stChatMessage"] table {
    border-collapse: collapse;
    width: 100%;
    background: #ffffff;
}
[data-testid="stChatMessage"] th {
    background: #2e7d32 !important;
    color: white !important;
    padding: 6px 10px;
}
[data-testid="stChatMessage"] td { padding: 5px 10px; border-bottom: 1px solid #e0f0e0; }
[data-testid="stChatMessage"] tr:nth-child(even) td { background: #f1f8f1; }

/* input de chat */
[data-testid="stChatInput"] textarea {
    background: #f9fdf9 !important;
    color: #1a2e1a !important;
    border: 1.5px solid #2e7d32 !important;
    border-radius: 10px !important;
}

/* botГµes primГЎrios */
.stButton > button {
    background: linear-gradient(135deg, #2e7d32, #1b5e20);
    color: white;
    border: none;
    border-radius: 8px;
    font-weight: 600;
    padding: 0.5rem 1.2rem;
    transition: all .2s;
}
.stButton > button:hover {
    background: linear-gradient(135deg, #388e3c, #2e7d32);
    transform: translateY(-1px);
}

/* download buttons */
a.download-btn {
    display: inline-block;
    background: linear-gradient(135deg, #1565c0, #0d47a1);
    color: white !important;
    border-radius: 8px;
    padding: 0.45rem 1.1rem;
    text-decoration: none;
    font-size: 0.85rem;
    font-weight: 600;
    margin-right: 8px;
    margin-top: 6px;
}
a.download-btn:hover { background: linear-gradient(135deg, #1976d2, #1565c0); }
a.download-btn.green { background: linear-gradient(135deg, #2e7d32, #1b5e20); }

/* tГ­tulo principal */
h1 { color: #2e7d32 !important; font-weight: 900; }
h2, h3 { color: #388e3c !important; }

/* file uploader */
[data-testid="stFileUploader"] {
    background: #f9fdf9;
    border: 1.5px dashed #2e7d32;
    border-radius: 10px;
}

/* divider */
hr { border-color: #c8e6c9 !important; }

/* expander */
details { background: #f9fdf9; border-radius: 8px; border: 1px solid #c8e6c9; }

/* disclaimer box */
.disclaimer-box {
    background: #fff8e1;
    border: 1.5px solid #f9a825;
    border-radius: 8px;
    padding: 10px 14px;
    color: #5d4037 !important;
    font-size: 0.88rem;
    margin-top: 12px;
}
</style>
""", unsafe_allow_html=True)

# в•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђ
# SYSTEM PROMPT
# в•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђ
SYSTEM_PROMPT = """VocГЄ Г© o **AgrГґnomo IA**, assistente tГ©cnico agrГ­cola especializado da Guara Agro.
Foi treinado com as principais referГЄncias brasileiras de fertilidade do solo: 5ВЄ AproximaГ§ГЈo CFSEMG (1999), EMBRAPA, Apostila de NutriГ§ГЈo de Plantas (UFV/UNIPAM).

**PropГіsito:** ajudar produtores rurais, agrГґnomos e tГ©cnicos a interpretar anГЎlises de solo, calcular calagem e adubaГ§ГЈo, diagnosticar deficiГЄncias nutricionais e recomendar manejos culturais.

## REGRAS DE COMPORTAMENTO
- Sempre em portuguГЄs do Brasil, linguagem clara para o produtor rural
- Responda de forma direta e objetiva
- Use emojis com moderaГ§ГЈo (рџЊ± вњ… вљ пёЏ)
- Nunca invente dados вЂ” se nГЈo sabe, diz claramente
- Este Г© um agente WEB вЂ” pode usar markdown completo (tabelas, negrito, listas)

## QUANDO FALTAM DADOS
Antes de calagem ou adubaГ§ГЈo, pergunte:
1. Textura do solo (% argila ou tipo: arenoso/mГ©dio/argiloso)
2. Cultura e produtividade esperada
3. Resultados da anГЎlise (pH, MO, P, K, Ca, Mg, Al, H+Al, CTC, V%)
4. PRNT do calcГЎrio
5. Sistema de plantio

Se o usuГЎrio enviar os dados de uma vez, calcule direto.

## SEQUГЉNCIA AO RECEBER LAUDO (foto ou texto)
1. Interprete cada parГўmetro (adequado / baixo / alto / excessivo)
2. Calcule NC pelos mГ©todos da 5ВЄ Aprox e apresente o mais conservador
3. Verifique necessidade de gessagem (AlВівЃє > 0,5 ou V% < 50 em subsolo)
4. Recomende NPK para a cultura/produtividade informada
5. Alerte sobre micronutrientes crГ­ticos (B, Zn, Cu, Mn)
6. Resumo final com tabela de aГ§Гµes prГЎticas

## FORMULADOS вЂ” REGRA MГљLTIPLAS OPГ‡Г•ES
Ao recomendar adubo, SEMPRE dГЄ 2 a 4 alternativas:

| Nutriente | OpГ§ГЈo 1 | OpГ§ГЈo 2 | OpГ§ГЈo 3 |
|-----------|---------|---------|---------|
| Pв‚‚Oв‚… | MAP 09-46-00 | 10-40-10 | TSP 00-46-00 |
| N | Ureia 45-00-00 | Sulfato amГґnio 21-00-00 | Nitrato amГґnio 33-00-00 |
| K | KCl 00-00-60 | Sulfato K 00-00-50 | KCl granulado 00-00-58 |

Formulados completos cerrado mineiro: 08-28-16, 10-20-20, 04-30-10, 05-25-15

## TABELAS
Sempre que tiver dados comparativos ou recomendaГ§Гµes, apresente em tabela markdown.
Exemplo de tabela de resultado:

| ParГўmetro | Resultado | ReferГЄncia | Status |
|-----------|-----------|------------|--------|
| pH Hв‚‚O | 5,1 | 5,5 вЂ“ 6,5 | вљ пёЏ Baixo |

## GERAГ‡ГѓO DE ARQUIVOS
Quando o usuГЎrio pedir PDF, Word, PowerPoint, relatГіrio ou planilha, NГѓO diga "clique no botГЈo" nem redirecione. Em vez disso, gere IMEDIATAMENTE o conteГєdo COMPLETO e DETALHADO do documento diretamente na sua resposta usando markdown rico:
- Use ## para seГ§Гµes principais
- Use ### para subseГ§Гµes
- Use tabelas markdown para todos os dados comparativos
- Use listas com bullets para recomendaГ§Гµes
- Inclua TODOS os cГЎlculos, valores e recomendaГ§Гµes

O sistema converterГЎ automaticamente sua resposta em PDF, Word, Excel e PowerPoint via os botГµes de download. Sua tarefa Г© gerar conteГєdo RICO e COMPLETO вЂ” nГЈo instruГ§Гµes sobre como usar o sistema.

## AVISO FINAL вЂ” SEMPRE INCLUIR
вљ пёЏ *Esta recomendaГ§ГЈo nГЈo substitui o laudo de Engenheiro AgrГґnomo habilitado (CREA).*
"""

# в•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђ
# AUTENTICAГ‡ГѓO
# в•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђ
VALID_CODES = st.secrets.get("ACCESS_CODES", "DEMO2024,GUARA2024,AGRO2024").split(",")

def check_auth():
    if st.session_state.get("authenticated"):
        return True
    st.markdown("## рџЊ± AgrГґnomo IA вЂ” Guara Agro")
    st.markdown("### Acesso ao Agente")
    code = st.text_input("CГіdigo de acesso:", type="password", placeholder="Digite seu cГіdigo...")
    if st.button("Entrar", use_container_width=True):
        if code.strip().upper() in [c.strip().upper() for c in VALID_CODES]:
            st.session_state.authenticated = True
            st.session_state.user_code = code.strip().upper()
            st.rerun()
        else:
            st.error("вќЊ CГіdigo invГЎlido. Verifique com a Guara Agro.")
    st.markdown("---")
    st.caption("Ainda nГЈo tem acesso? Entre em contato via WhatsApp.")
    return False

# в•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђ
# GERAГ‡ГѓO DE PDF
# в•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђ
def markdown_to_pdf(markdown_text: str, title: str = "RelatГіrio вЂ” AgrГґnomo IA") -> bytes:
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        rightMargin=2*cm, leftMargin=2*cm,
        topMargin=2.5*cm, bottomMargin=2*cm
    )

    styles = getSampleStyleSheet()
    verde_escuro = colors.HexColor("#1b5e20")
    verde_medio  = colors.HexColor("#2e7d32")
    verde_claro  = colors.HexColor("#e8f5e9")
    cinza_bg     = colors.HexColor("#f5f5f5")

    estilo_titulo = ParagraphStyle(
        "Titulo", parent=styles["Title"],
        textColor=verde_escuro, fontSize=20, spaceAfter=6,
        fontName="Helvetica-Bold"
    )
    estilo_h2 = ParagraphStyle(
        "H2", parent=styles["Heading2"],
        textColor=verde_medio, fontSize=14, spaceBefore=14, spaceAfter=4,
        fontName="Helvetica-Bold"
    )
    estilo_h3 = ParagraphStyle(
        "H3", parent=styles["Heading3"],
        textColor=verde_escuro, fontSize=12, spaceBefore=10, spaceAfter=3,
        fontName="Helvetica-Bold"
    )
    estilo_body = ParagraphStyle(
        "Body", parent=styles["Normal"],
        fontSize=10, leading=15, spaceAfter=4,
        fontName="Helvetica"
    )
    estilo_aviso = ParagraphStyle(
        "Aviso", parent=styles["Normal"],
        fontSize=9, textColor=colors.HexColor("#b71c1c"),
        borderPad=4, backColor=colors.HexColor("#fff3e0"),
        borderColor=colors.HexColor("#e65100"), borderWidth=0.5,
        fontName="Helvetica-Oblique", leading=13
    )

    story = []

    # CabeГ§alho
    story.append(Paragraph("рџЊ± GUARA AGRO", ParagraphStyle(
        "Logo", parent=styles["Normal"],
        textColor=verde_medio, fontSize=11, fontName="Helvetica-Bold"
    )))
    story.append(Spacer(1, 4))
    story.append(Paragraph(title, estilo_titulo))
    story.append(Paragraph(
        f"Gerado em {datetime.now().strftime('%d/%m/%Y Г s %H:%M')} В· AgrГґnomo IA v1.0",
        ParagraphStyle("sub", parent=styles["Normal"], fontSize=9,
                       textColor=colors.grey, fontName="Helvetica")
    ))
    story.append(HRFlowable(width="100%", thickness=2, color=verde_medio, spaceAfter=12))

    # Parser linha a linha
    lines = markdown_text.split("\n")
    i = 0
    while i < len(lines):
        line = lines[i]

        # Tabelas markdown
        if "|" in line and i + 1 < len(lines) and "|" in lines[i+1] and "---" in lines[i+1]:
            table_lines = []
            while i < len(lines) and "|" in lines[i]:
                cols = [c.strip() for c in lines[i].split("|") if c.strip()]
                if "---" not in lines[i]:
                    table_lines.append(cols)
                i += 1
            if table_lines:
                max_cols = max(len(r) for r in table_lines)
                data = [r + [""] * (max_cols - len(r)) for r in table_lines]
                col_width = (A4[0] - 4*cm) / max_cols
                t = Table(data, colWidths=[col_width] * max_cols)
                t.setStyle(TableStyle([
                    ("BACKGROUND",  (0, 0), (-1, 0), verde_medio),
                    ("TEXTCOLOR",   (0, 0), (-1, 0), colors.white),
                    ("FONTNAME",    (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE",    (0, 0), (-1, -1), 9),
                    ("BACKGROUND",  (0, 1), (-1, -1), cinza_bg),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, cinza_bg]),
                    ("GRID",        (0, 0), (-1, -1), 0.5, colors.HexColor("#c8e6c9")),
                    ("ALIGN",       (0, 0), (-1, -1), "LEFT"),
                    ("VALIGN",      (0, 0), (-1, -1), "MIDDLE"),
                    ("PADDING",     (0, 0), (-1, -1), 5),
                    ("TOPPADDING",  (0, 0), (-1, -1), 4),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                ]))
                story.append(t)
                story.append(Spacer(1, 8))
            continue

        # TГ­tulos
        if line.startswith("### "):
            story.append(Paragraph(line[4:].strip(), estilo_h3))
        elif line.startswith("## "):
            story.append(Paragraph(line[3:].strip(), estilo_h2))
        elif line.startswith("# "):
            story.append(Paragraph(line[2:].strip(), estilo_titulo))

        # Listas
        elif line.strip().startswith(("- ", "* ", "вЂў ")):
            txt = line.strip()[2:].strip()
            txt = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', txt)
            story.append(Paragraph(f"вЂў {txt}", ParagraphStyle(
                "li", parent=estilo_body, leftIndent=14, spaceAfter=2
            )))
        elif re.match(r'^\d+\.', line.strip()):
            txt = line.strip()
            txt = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', txt)
            story.append(Paragraph(txt, ParagraphStyle(
                "ol", parent=estilo_body, leftIndent=14, spaceAfter=2
            )))

        # Aviso вљ пёЏ
        elif "вљ пёЏ" in line or "nГЈo substitui" in line.lower():
            clean = re.sub(r'\*(.+?)\*', r'<i>\1</i>', line.strip().lstrip("вљ пёЏ *"))
            story.append(Spacer(1, 6))
            story.append(Paragraph(f"вљ пёЏ {clean}", estilo_aviso))

        # HR
        elif line.strip() in ("---", "___", "***"):
            story.append(HRFlowable(width="100%", thickness=0.5,
                                     color=colors.HexColor("#c8e6c9"), spaceAfter=6))

        # Linha vazia
        elif line.strip() == "":
            story.append(Spacer(1, 4))

        # ParГЎgrafo normal
        else:
            txt = line.strip()
            txt = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', txt)
            txt = re.sub(r'\*(.+?)\*',   r'<i>\1</i>', txt)
            txt = re.sub(r'`(.+?)`',     r'<font name="Courier">\1</font>', txt)
            if txt:
                story.append(Paragraph(txt, estilo_body))

        i += 1

    # RodapГ©
    story.append(Spacer(1, 16))
    story.append(HRFlowable(width="100%", thickness=1, color=verde_medio))
    story.append(Spacer(1, 4))
    story.append(Paragraph(
        "Guara Agro В· AgrГґnomo IA v1.0 В· вљ пёЏ Esta recomendaГ§ГЈo nГЈo substitui o laudo de Engenheiro AgrГґnomo habilitado (CREA).",
        ParagraphStyle("footer", parent=styles["Normal"], fontSize=8,
                       textColor=colors.grey, alignment=1, fontName="Helvetica-Oblique")
    ))

    doc.build(story)
    buffer.seek(0)
    return buffer.read()


# в•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђ
# GERAГ‡ГѓO DE EXCEL
# в•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђ
def extract_tables_from_markdown(text: str):
    """Extrai tabelas markdown e retorna lista de (header, rows)."""
    tables = []
    lines = text.split("\n")
    i = 0
    while i < len(lines):
        if "|" in lines[i]:
            raw = []
            while i < len(lines) and "|" in lines[i]:
                cols = [c.strip() for c in lines[i].split("|") if c.strip()]
                if "---" not in lines[i]:
                    raw.append(cols)
                i += 1
            if len(raw) >= 2:
                tables.append((raw[0], raw[1:]))
        else:
            i += 1
    return tables


def response_to_excel(markdown_text: str, title: str = "AgrГґnomo IA") -> bytes:
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "RecomendaГ§Гµes"

    verde = "1B5E20"
    verde_claro = "E8F5E9"
    verde_medio = "2E7D32"
    amarelo = "FFF9C4"

    # CabeГ§alho
    ws.merge_cells("A1:G1")
    c = ws["A1"]
    c.value = f"рџЊ± GUARA AGRO вЂ” {title}"
    c.font = Font(bold=True, size=16, color="FFFFFF")
    c.fill = PatternFill("solid", fgColor=verde_medio)
    c.alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 32

    ws.merge_cells("A2:G2")
    c2 = ws["A2"]
    c2.value = f"Gerado em {datetime.now().strftime('%d/%m/%Y %H:%M')} В· AgrГґnomo IA v1.0"
    c2.font = Font(italic=True, size=9, color="666666")
    c2.alignment = Alignment(horizontal="center")

    row = 4

    tables = extract_tables_from_markdown(markdown_text)

    if tables:
        for header, rows in tables:
            # Header da tabela
            for col_idx, h in enumerate(header, 1):
                cell = ws.cell(row=row, column=col_idx, value=h)
                cell.font = Font(bold=True, color="FFFFFF", size=11)
                cell.fill = PatternFill("solid", fgColor=verde)
                cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
                cell.border = Border(
                    bottom=Side(style="medium", color=verde_medio),
                    right=Side(style="thin", color="CCCCCC")
                )
                ws.column_dimensions[cell.column_letter].width = max(20, len(str(h)) + 4)
            ws.row_dimensions[row].height = 22
            row += 1

            # Linhas
            for r_idx, data_row in enumerate(rows):
                bg = verde_claro if r_idx % 2 == 0 else "FFFFFF"
                for col_idx, val in enumerate(data_row, 1):
                    cell = ws.cell(row=row, column=col_idx, value=val)
                    cell.fill = PatternFill("solid", fgColor=bg)
                    cell.alignment = Alignment(vertical="center", wrap_text=True)
                    cell.border = Border(
                        bottom=Side(style="thin", color="E0E0E0"),
                        right=Side(style="thin", color="E0E0E0")
                    )
                    cell.font = Font(size=10)
                ws.row_dimensions[row].height = 18
                row += 1
            row += 1

    else:
        # Sem tabela вЂ” dump do texto
        linhas = [l for l in markdown_text.split("\n") if l.strip()]
        ws.cell(row=row, column=1, value="AnГЎlise / RecomendaГ§ГЈo").font = Font(bold=True)
        row += 1
        for linha in linhas:
            ws.cell(row=row, column=1, value=linha)
            row += 1

    # RodapГ©
    row += 1
    ws.merge_cells(f"A{row}:G{row}")
    c_footer = ws[f"A{row}"]
    c_footer.value = "вљ пёЏ Esta recomendaГ§ГЈo nГЈo substitui o laudo de Engenheiro AgrГґnomo habilitado (CREA). В· Guara Agro"
    c_footer.font = Font(italic=True, size=8, color="999999")
    c_footer.alignment = Alignment(horizontal="center")

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.read()


# в•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђ
# GERAГ‡ГѓO DE WORD (.docx)
# в•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђ
def markdown_to_docx(markdown_text: str, title: str = "RelatГіrio вЂ” AgrГґnomo IA") -> bytes:
    doc = Document()

    VERDE = RGBColor(0x2E, 0x7D, 0x32)
    VERDE_ESC = RGBColor(0x1B, 0x5E, 0x20)
    CINZA = RGBColor(0x75, 0x75, 0x75)

    # Margens
    for section in doc.sections:
        section.top_margin    = Cm(2)
        section.bottom_margin = Cm(2)
        section.left_margin   = Cm(2.5)
        section.right_margin  = Cm(2.5)

    # CabeГ§alho
    hdr = doc.add_paragraph()
    hdr.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = hdr.add_run("рџЊ± GUARA AGRO вЂ” AGRГ”NOMO IA")
    run.font.bold = True
    run.font.size = Pt(18)
    run.font.color.rgb = VERDE_ESC

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = sub.add_run(f"{title}  В·  {datetime.now().strftime('%d/%m/%Y %H:%M')}")
    r2.font.size = Pt(10)
    r2.font.color.rgb = CINZA
    r2.font.italic = True

    doc.add_paragraph()

    # Processa markdown linha a linha
    lines = markdown_text.split("\n")
    i = 0
    while i < len(lines):
        line = lines[i]

        # Headings
        if line.startswith("### "):
            p = doc.add_heading(line[4:].strip(), level=3)
            p.runs[0].font.color.rgb = VERDE
        elif line.startswith("## "):
            p = doc.add_heading(line[3:].strip(), level=2)
            p.runs[0].font.color.rgb = VERDE_ESC
        elif line.startswith("# "):
            p = doc.add_heading(line[2:].strip(), level=1)
            p.runs[0].font.color.rgb = VERDE_ESC

        # Tabela markdown
        elif line.strip().startswith("|") and i + 1 < len(lines) and "---" in lines[i + 1]:
            headers = [c.strip() for c in line.split("|") if c.strip()]
            i += 2  # pula linha de separador
            rows_data = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                row = [c.strip() for c in lines[i].split("|") if c.strip()]
                if row:
                    rows_data.append(row)
                i += 1

            n_cols = max(len(headers), max((len(r) for r in rows_data), default=0))
            table = doc.add_table(rows=1 + len(rows_data), cols=n_cols)
            table.style = "Table Grid"

            # Header row
            for j, h in enumerate(headers):
                cell = table.rows[0].cells[j]
                cell.text = h
                run = cell.paragraphs[0].runs[0]
                run.bold = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                shd = OxmlElement("w:shd")
                shd.set(qn("w:fill"), "2E7D32")
                shd.set(qn("w:color"), "auto")
                shd.set(qn("w:val"), "clear")
                tcPr.append(shd)

            # Data rows
            for ri, row_data in enumerate(rows_data):
                for ci, val in enumerate(row_data):
                    if ci < n_cols:
                        table.rows[ri + 1].cells[ci].text = val

            doc.add_paragraph()
            continue  # jГЎ avanГ§ou i dentro do loop

        # Bullet
        elif line.startswith("- ") or line.startswith("* "):
            txt = line[2:].strip()
            txt = re.sub(r'\*\*(.+?)\*\*', r'\1', txt)
            p = doc.add_paragraph(txt, style="List Bullet")

        # Linha vazia
        elif line.strip() == "":
            doc.add_paragraph()

        # ParГЎgrafo normal
        else:
            txt = line.strip()
            if not txt:
                i += 1
                continue
            p = doc.add_paragraph()
            # bold inline
            parts = re.split(r'\*\*(.+?)\*\*', txt)
            for pi, part in enumerate(parts):
                r = p.add_run(part)
                r.bold = (pi % 2 == 1)

        i += 1

    # RodapГ©
    doc.add_paragraph()
    footer_p = doc.add_paragraph()
    footer_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    fr = footer_p.add_run("вљ пёЏ Esta recomendaГ§ГЈo nГЈo substitui o laudo de Engenheiro AgrГґnomo habilitado (CREA).  В·  Guara Agro")
    fr.font.size = Pt(8)
    fr.font.italic = True
    fr.font.color.rgb = CINZA

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


# в•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђ
# GERAГ‡ГѓO DE POWERPOINT (.pptx)
# в•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђ
def markdown_to_pptx(markdown_text: str, title: str = "AnГЎlise AgrГґnomo IA") -> bytes:
    prs = Presentation()
    prs.slide_width  = PInches(13.33)
    prs.slide_height = PInches(7.5)

    VERDE     = PRGBColor(0x2E, 0x7D, 0x32)
    VERDE_ESC = PRGBColor(0x1B, 0x5E, 0x20)
    BRANCO    = PRGBColor(0xFF, 0xFF, 0xFF)
    CINZA_ESC = PRGBColor(0x33, 0x33, 0x33)

    blank_layout = prs.slide_layouts[6]  # totalmente em branco

    def add_rect(slide, x, y, w, h, fill_color):
        shape = slide.shapes.add_shape(1, PInches(x), PInches(y), PInches(w), PInches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    def add_textbox(slide, text, x, y, w, h, size=18, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
        txb = slide.shapes.add_textbox(PInches(x), PInches(y), PInches(w), PInches(h))
        tf  = txb.text_frame
        tf.word_wrap = True
        p   = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size  = PPt(size)
        run.font.bold  = bold
        run.font.italic = italic
        if color:
            run.font.color.rgb = color
        return txb

    # в”Ђв”Ђ Slide 1: Capa в”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђ
    slide1 = prs.slides.add_slide(blank_layout)
    add_rect(slide1, 0, 0, 13.33, 7.5, VERDE_ESC)
    add_rect(slide1, 0, 5.8, 13.33, 1.7, PRGBColor(0x14, 0x46, 0x14))
    add_textbox(slide1, "рџЊ± GUARA AGRO", 0.5, 1.5, 12.33, 1.2, size=32, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    add_textbox(slide1, "AGRГ”NOMO IA", 0.5, 2.5, 12.33, 1.0, size=44, bold=True, color=PRGBColor(0xA5, 0xD6, 0xA7), align=PP_ALIGN.CENTER)
    add_textbox(slide1, title, 0.5, 3.6, 12.33, 0.7, size=20, color=BRANCO, align=PP_ALIGN.CENTER)
    add_textbox(slide1, datetime.now().strftime("%d/%m/%Y"), 0.5, 6.1, 12.33, 0.6, size=14, color=PRGBColor(0xA5, 0xD6, 0xA7), align=PP_ALIGN.CENTER)

    # в”Ђв”Ђ Processa conteГєdo em slides в”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђ
    lines = markdown_text.split("\n")
    slides_data = []  # lista de (titulo, [bullets])
    current_title = title
    current_bullets = []

    for line in lines:
        if line.startswith("## ") or line.startswith("# "):
            if current_bullets:
                slides_data.append((current_title, current_bullets))
                current_bullets = []
            current_title = re.sub(r'[#*]+', '', line).strip()
        elif line.startswith("### "):
            current_bullets.append(("sub", re.sub(r'#+', '', line).strip()))
        elif line.startswith("- ") or line.startswith("* "):
            txt = re.sub(r'\*\*(.+?)\*\*', r'\1', line[2:].strip())
            current_bullets.append(("bullet", txt))
        elif line.strip() and not line.startswith("|") and "---" not in line:
            txt = re.sub(r'\*\*(.+?)\*\*', r'\1', line.strip())
            current_bullets.append(("text", txt))

    if current_bullets:
        slides_data.append((current_title, current_bullets))

    # в”Ђв”Ђ Gera slides de conteГєdo в”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђв”Ђ
    for slide_title, bullets in slides_data:
        slide = prs.slides.add_slide(blank_layout)

        # Barra superior verde
        add_rect(slide, 0, 0, 13.33, 1.2, VERDE)
        add_textbox(slide, slide_title, 0.3, 0.1, 12.73, 1.0, size=24, bold=True, color=BRANCO)

        # Barra inferior
        add_rect(slide, 0, 7.1, 13.33, 0.4, VERDE_ESC)
        add_textbox(slide, "Guara Agro В· AgrГґnomo IA  |  вљ пёЏ NГЈo substitui laudo de Engenheiro AgrГґnomo (CREA)", 0.2, 7.1, 12.9, 0.38, size=9, color=BRANCO, italic=True)

        # ConteГєdo
        y = 1.35
        max_y = 6.9
        for kind, txt in bullets:
            if y >= max_y:
                break
            if kind == "sub":
                add_textbox(slide, txt, 0.4, y, 12.5, 0.45, size=14, bold=True, color=VERDE)
                y += 0.5
            elif kind == "bullet":
                add_textbox(slide, f"вЂў {txt}", 0.6, y, 12.1, 0.4, size=13, color=CINZA_ESC)
                y += 0.42
            else:
                add_textbox(slide, txt, 0.4, y, 12.5, 0.4, size=12, color=CINZA_ESC)
                y += 0.45

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# в•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђ
# HELPER: download link HTML
# в•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђ
def make_download_link(data: bytes, filename: str, label: str, css_class: str = "") -> str:
    b64 = base64.b64encode(data).decode()
    ext = filename.rsplit(".", 1)[-1].lower()
    mimes = {
        "pdf":  "application/pdf",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }
    mime = mimes.get(ext, "application/octet-stream")
    return (
        f'<a href="data:{mime};base64,{b64}" download="{filename}" '
        f'class="download-btn {css_class}">{label}</a>'
    )


# в•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђ
# CLAUDE API
# в•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђ
def get_client():
    api_key = st.secrets.get("ANTHROPIC_API_KEY", "")
    if not api_key:
        st.error("ANTHROPIC_API_KEY nГЈo configurada em .streamlit/secrets.toml")
        st.stop()
    return anthropic.Anthropic(api_key=api_key)


def call_claude(messages: list, image_b64: str = None, image_mime: str = None) -> str:
    client = get_client()

    api_messages = []
    for m in messages[:-1]:  # histГіrico
        api_messages.append({"role": m["role"], "content": m["content"]})

    # Гљltima mensagem (pode ter imagem)
    last = messages[-1]
    if image_b64 and last["role"] == "user":
        content = [
            {"type": "image", "source": {"type": "base64", "media_type": image_mime, "data": image_b64}},
            {"type": "text", "text": last["content"]},
        ]
    else:
        content = last["content"]

    api_messages.append({"role": "user", "content": content})

    with st.spinner("рџЊ± Analisando..."):
        response = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=4096,
            system=SYSTEM_PROMPT,
            messages=api_messages,
        )
    return response.content[0].text


# в•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђ
# INTERFACE PRINCIPAL
# в•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђ
def main():
    # Header
    col1, col2 = st.columns([5, 2])
    with col1:
        st.markdown("# рџЊ± AgrГґnomo IA")
        st.caption("Guara Agro В· Assistente tГ©cnico especializado em solo e nutriГ§ГЈo")
    with col2:
        st.markdown("<br>", unsafe_allow_html=True)
        if st.button("рџљЄ Sair", use_container_width=True):
            st.session_state.clear()
            st.rerun()

    st.divider()

    # Inicializa histГіrico
    if "messages" not in st.session_state:
        st.session_state.messages = []
        st.session_state.messages.append({
            "role": "assistant",
            "content": (
                "рџ‘‹ OlГЎ! Sou o **AgrГґnomo IA** da Guara Agro.\n\n"
                "Posso te ajudar com:\n"
                "- рџ“Љ **InterpretaГ§ГЈo de laudos de solo**\n"
                "- рџ§® **CГЎlculo de calagem e adubaГ§ГЈo NPK**\n"
                "- рџ”Ќ **DiagnГіstico de deficiГЄncias nutricionais**\n"
                "- рџЊї **RecomendaГ§Гµes para o Cerrado Mineiro**\n\n"
                "Envie o laudo (foto ou texto), faГ§a uma pergunta ou me peГ§a para gerar um relatГіrio. "
                "Depois de cada anГЎlise, vocГЄ poderГЎ **baixar o relatГіrio em PDF ou Excel** direto aqui.\n\n"
                "вљ пёЏ *Esta recomendaГ§ГЈo nГЈo substitui o laudo de Engenheiro AgrГґnomo habilitado (CREA).*"
            ),
        })

    # Upload de imagem (sidebar)
    with st.sidebar:
        st.markdown("### рџ“Ћ Enviar Laudo / Foto")
        uploaded = st.file_uploader(
            "Foto do laudo, planta ou solo",
            type=["jpg", "jpeg", "png", "webp"],
            help="Envie foto do laudo de solo, planta com deficiГЄncia ou ГЎrea problema"
        )
        if uploaded:
            st.image(uploaded, caption="Imagem carregada", use_container_width=True)

        st.divider()
        st.markdown("### рџ’Ў Exemplos de perguntas")
        exemplos = [
            "Interprete este laudo de solo",
            "Calcule calagem para pH 6,0",
            "Recomende NPK para milho 120 sc/ha",
            "O que causa clorose nas folhas?",
            "Qual formulado para soja no cerrado?",
            "Gere um relatГіrio em PDF",
        ]
        for ex in exemplos:
            if st.button(ex, use_container_width=True, key=f"ex_{ex}"):
                st.session_state.exemplo_selecionado = ex
                st.rerun()

        st.divider()
        if st.button("рџ—‘пёЏ Limpar conversa", use_container_width=True):
            st.session_state.messages = []
            st.rerun()

    DISCLAIMER = "вљ пёЏ **Esta recomendaГ§ГЈo nГЈo substitui o laudo de Engenheiro AgrГґnomo habilitado (CREA).**"

    def exibe_msg_assistente(content: str, show_disclaimer: bool = True):
        # Remove disclaimer inline que o modelo possa ter incluГ­do
        texto = re.sub(r'вљ пёЏ\s*\*?Esta recomendaГ§ГЈo nГЈo substitui.*?CREA\.\*?\n?', '', content, flags=re.IGNORECASE).strip()
        st.markdown(texto)
        if show_disclaimer:
            st.markdown(
                f'<div class="disclaimer-box">{DISCLAIMER}</div>',
                unsafe_allow_html=True
            )

    # Exibe histГіrico
    for i, msg in enumerate(st.session_state.messages):
        avatar = "рџЊ±" if msg["role"] == "assistant" else "рџ‘¤"
        with st.chat_message(msg["role"], avatar=avatar):
            if msg["role"] == "assistant":
                # Mensagem de boas-vindas (Г­ndice 0) mostra disclaimer simples
                exibe_msg_assistente(msg["content"], show_disclaimer=(i > 0))
            else:
                st.markdown(msg["content"])

    # BotГµes de download da Гєltima resposta do assistente
    assistant_msgs = [m for m in st.session_state.messages if m["role"] == "assistant"]
    if len(assistant_msgs) > 1:  # exclui mensagem de boas-vindas
        last_ai = assistant_msgs[-1]["content"]
        nome_data = datetime.now().strftime("%Y%m%d_%H%M")

        pdf_bytes   = markdown_to_pdf(last_ai, "AnГЎlise AgrГґnomo IA")
        excel_bytes = response_to_excel(last_ai, "AnГЎlise AgrГґnomo IA")
        docx_bytes  = markdown_to_docx(last_ai, "AnГЎlise AgrГґnomo IA")
        pptx_bytes  = markdown_to_pptx(last_ai, "AnГЎlise AgrГґnomo IA")

        link_pdf   = make_download_link(pdf_bytes,   f"agronomo_ia_{nome_data}.pdf",  "рџ“Ґ PDF",   "")
        link_excel = make_download_link(excel_bytes, f"agronomo_ia_{nome_data}.xlsx", "рџ“Љ Excel", "green")
        link_docx  = make_download_link(docx_bytes,  f"agronomo_ia_{nome_data}.docx", "рџ“ќ Word",  "")
        link_pptx  = make_download_link(pptx_bytes,  f"agronomo_ia_{nome_data}.pptx", "рџ“‘ PPT",   "green")

        st.markdown(
            f'<div style="margin-bottom:8px">{link_pdf} {link_excel} {link_docx} {link_pptx}</div>',
            unsafe_allow_html=True
        )

    # Input
    image_b64, image_mime = None, None

    # Exemplo selecionado na sidebar
    prompt = st.session_state.pop("exemplo_selecionado", None)

    # Chat input principal
    chat_input = st.chat_input(
        "Envie o laudo, faГ§a uma pergunta ou peГ§a um relatГіrio...",
        key="chat_input"
    )
    if chat_input:
        prompt = chat_input

    if prompt:
        # Prepara imagem se houver
        if uploaded:
            img_bytes = uploaded.read()
            image_b64  = base64.b64encode(img_bytes).decode()
            image_mime = uploaded.type
            display_content = f"рџ“· *[Imagem enviada: {uploaded.name}]*\n\n{prompt}"
        else:
            display_content = prompt

        # Mostra mensagem do usuГЎrio
        with st.chat_message("user", avatar="рџ‘¤"):
            st.markdown(display_content)

        st.session_state.messages.append({"role": "user", "content": display_content})

        # Chama Claude
        resposta = call_claude(st.session_state.messages, image_b64, image_mime)

        # Mostra resposta
        with st.chat_message("assistant", avatar="рџЊ±"):
            exibe_msg_assistente(resposta)

        st.session_state.messages.append({"role": "assistant", "content": resposta})
        st.rerun()


# в•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђ
# ENTRY POINT
# в•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђв•ђ
if not check_auth():
    st.stop()

main()
